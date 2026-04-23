"""Generate Phase 1 Roadway Data Pipeline PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Color palette (Jacobs-inspired professional blues) --
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MED_BLUE = RGBColor(0x00, 0x5E, 0xA2)
LIGHT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TABLE_HEADER_BG = RGBColor(0x00, 0x2B, 0x5C)
TABLE_ALT_ROW = RGBColor(0xE8, 0xF4, 0xFC)


def add_bottom_bar(slide):
    """Add a thin accent bar at the bottom of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), prs.slide_height - Inches(0.15),
        prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()


def add_slide_number(slide, num, total):
    txBox = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.55),
        Inches(1), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.RIGHT


def set_cell_text(cell, text, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_table_header(row, color=WHITE, bg=TABLE_HEADER_BG, font_size=12):
    for cell in row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


TOTAL_SLIDES = 11

# =========================================================
# SLIDE 1: Title Slide
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Accent stripe
stripe = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(3.4),
    prs.slide_width, Inches(0.08)
)
stripe.fill.solid()
stripe.fill.fore_color.rgb = ACCENT_TEAL
stripe.line.fill.background()

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Georgia Statewide Roadway Data Pipeline"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = "Phase 1 — Foundation Layer"
p2.font.size = Pt(28)
p2.font.color.rgb = ACCENT_TEAL
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(8)

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.text = "GDOT RAPTOR Pipeline — Status Briefing"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
p3.alignment = PP_ALIGN.LEFT

p4 = tf2.add_paragraph()
p4.text = "April 2026"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x88, 0xAA, 0xCC)
p4.alignment = PP_ALIGN.LEFT
p4.space_before = Pt(12)

# =========================================================
# SLIDE 2: Executive Summary
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Executive Summary"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# Separator
sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Summary text
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True

bullets = [
    "Phase 1 builds the foundational roadway layer that RAPTOR scoring runs on top of.",
    "Combines 8 GDOT data sources into a single statewide roadway network with traffic and roadway attributes.",
    "Produces a staged database and spatial file ready for downstream RAPTOR analysis.",
    "Status: Complete for current project scope — 116 / 116 validation checks passing.",
]
for i, txt in enumerate(bullets):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = txt
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(14)
    p.level = 0

# Key numbers box (right side)
box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.5), Inches(1.5), Inches(4.2), Inches(4.8)
)
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
box.line.color.rgb = MED_BLUE
box.line.width = Pt(1.5)

stats = [
    ("245,863", "Roadway Segments"),
    ("99.96%", "AADT Coverage"),
    ("8", "Data Sources Combined"),
    ("118", "Database Columns"),
    ("159", "Counties Statewide"),
]

for i, (val, label) in enumerate(stats):
    y = Inches(1.7) + Inches(i * 0.9)
    vBox = slide.shapes.add_textbox(Inches(8.8), y, Inches(3.6), Inches(0.45))
    vtf = vBox.text_frame
    vp = vtf.paragraphs[0]
    vp.text = val
    vp.font.size = Pt(26)
    vp.font.bold = True
    vp.font.color.rgb = MED_BLUE
    vp.alignment = PP_ALIGN.CENTER

    lBox = slide.shapes.add_textbox(Inches(8.8), y + Inches(0.4), Inches(3.6), Inches(0.3))
    ltf = lBox.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(12)
    lp.font.color.rgb = MED_GRAY
    lp.alignment = PP_ALIGN.CENTER


# =========================================================
# SLIDE 3: Data Sources
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Data Sources"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

sources = [
    ("#", "Source", "What We Get"),
    ("1", "GDOT Road Inventory GDB", "Base road network — route geometry, route IDs, milepoints, roadway attributes (lanes, surface, functional class, median, shoulders)"),
    ("2", "GDOT Traffic GDB", "Current-year AADT, truck AADT, VMT, K-factor, D-factor — traffic interval boundaries drive segmentation"),
    ("3", "FHWA HPMS (2020–2024)", "Parallel GDOT-official AADT, pavement condition (IRI, rutting, cracking), signed-route classification"),
    ("4", "GDOT GPAS SpeedZone (On)", "Posted speed limits for state highways — matched by route ID + milepoint"),
    ("5", "GDOT GPAS SpeedZone (Off)", "Posted speed limits for local roads — matched by road name + county"),
    ("6", "GDOT GPAS Reference Layers", "Authoritative signed-route verification (Interstate / US / State Route)"),
    ("7", "GDOT Boundaries Service", "County (159) and district (7) boundary polygons — split-driving: routes segmented at every boundary crossing"),
    ("8", "GDOT EOC Evacuation Routes", "Hurricane evacuation route flags — secondary split at corridor boundaries"),
]

rows, cols = len(sources), 3
col_widths = [Inches(0.5), Inches(3.2), Inches(8.0)]
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), sum(col_widths), Inches(5.5)).table
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = w

for ri, row_data in enumerate(sources):
    for ci, val in enumerate(row_data):
        set_cell_text(tbl.cell(ri, ci), val, font_size=11,
                      bold=(ri == 0),
                      color=WHITE if ri == 0 else DARK_GRAY,
                      alignment=PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT)
    if ri == 0:
        style_table_header(tbl.rows[ri])
    elif ri % 2 == 0:
        for ci in range(cols):
            tbl.cell(ri, ci).fill.solid()
            tbl.cell(ri, ci).fill.fore_color.rgb = TABLE_ALT_ROW


# =========================================================
# SLIDE 4: Pipeline Flow (overview)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Pipeline Flow — 12-Step Process"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

steps_left = [
    ("1", "Load official road network", "206,994 routes from GDOT GDB"),
    ("2", "Attach roadway attributes", "15 attribute layers joined by route ID"),
    ("3", "Segment: traffic + geography", "Split at traffic intervals AND admin boundaries"),
    ("4", "Post-split overlay flags", "Legislative districts, city (non-splitting)"),
    ("5", "Parse route identity", "ROUTE_ID → family, number, suffix"),
    ("6", "Enrich with speed limits", "GPAS on-system + off-system matching"),
]

steps_right = [
    ("7", "Add HPMS data", "Parallel GDOT-official AADT + pavement"),
    ("8", "Verify signed-route class", "GPAS authoritative override"),
    ("9", "Project future AADT 2044", "4-step fill chain → 99.96% coverage"),
    ("10", "Flag evacuation routes", "EOC corridor overlay + secondary split"),
    ("11", "Derive RAPTOR fields", "PCT_SADT, PCT_CADT, HWY_DES"),
    ("12", "Write outputs", "SQLite DB + GeoPackage + CSV"),
]

for col_idx, steps in enumerate([steps_left, steps_right]):
    x_base = Inches(0.8) + Inches(col_idx * 6.2)
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.5) + Inches(i * 0.95)

        # Number circle
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x_base, y, Inches(0.45), Inches(0.45)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = MED_BLUE
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.paragraphs[0].text = num
        ctf.paragraphs[0].font.size = Pt(14)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        tBox = slide.shapes.add_textbox(x_base + Inches(0.6), y - Inches(0.02), Inches(5), Inches(0.35))
        tp = tBox.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(15)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BLUE

        # Description
        dBox = slide.shapes.add_textbox(x_base + Inches(0.6), y + Inches(0.3), Inches(5), Inches(0.35))
        dp = dBox.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = MED_GRAY


# =========================================================
# SLIDE 5: Segmentation Deep Dive
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Dual-Source Segmentation"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Subtitle
txBox_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11), Inches(0.5))
tf_sub = txBox_sub.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Routes are split where traffic values change AND where they cross administrative boundaries"
p_sub.font.size = Pt(15)
p_sub.font.color.rgb = MED_GRAY
p_sub.font.italic = True

# Left box — Traffic breakpoints
traf_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.8)
)
traf_box.fill.solid()
traf_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
traf_box.line.color.rgb = MED_BLUE
traf_box.line.width = Pt(2)

traf_title = slide.shapes.add_textbox(Inches(1.1), Inches(2.15), Inches(5.0), Inches(0.5))
ttp = traf_title.text_frame.paragraphs[0]
ttp.text = "Traffic Interval Breakpoints"
ttp.font.size = Pt(18)
ttp.font.bold = True
ttp.font.color.rgb = MED_BLUE

traf_items = [
    "FROM/TO milepoints from GDOT Traffic GDB",
    "Split wherever AADT, truck AADT, or VMT changes",
    "46,029 traffic records define interval boundaries",
    "Each segment gets one consistent traffic record",
]
for i, txt in enumerate(traf_items):
    y = Inches(2.7) + Inches(i * 0.45)
    bBox = slide.shapes.add_textbox(Inches(1.3), y, Inches(4.8), Inches(0.4))
    btf = bBox.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = f"\u2022  {txt}"
    bp.font.size = Pt(13)
    bp.font.color.rgb = DARK_GRAY

# Right box — Admin boundary breakpoints
admin_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.8)
)
admin_box.fill.solid()
admin_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
admin_box.line.color.rgb = RGBColor(0x4C, 0xAF, 0x50)
admin_box.line.width = Pt(2)

admin_title = slide.shapes.add_textbox(Inches(7.3), Inches(2.15), Inches(5.2), Inches(0.5))
atp = admin_title.text_frame.paragraphs[0]
atp.text = "Administrative Boundary Breakpoints"
atp.font.size = Pt(18)
atp.font.bold = True
atp.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

admin_items = [
    "County boundaries (159 counties)",
    "GDOT District boundaries (7 districts)",
    "Area Office boundaries",
    "MPO boundaries",
    "Regional Commission boundaries",
]
for i, txt in enumerate(admin_items):
    y = Inches(2.7) + Inches(i * 0.4)
    bBox = slide.shapes.add_textbox(Inches(7.5), y, Inches(5.0), Inches(0.4))
    btf = bBox.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = f"\u2022  {txt}"
    bp.font.size = Pt(13)
    bp.font.color.rgb = DARK_GRAY

# Center merge arrow
merge_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.8), Inches(5.1), Inches(5.7), Inches(0.6)
)
merge_box.fill.solid()
merge_box.fill.fore_color.rgb = DARK_BLUE
merge_box.line.fill.background()
mtf = merge_box.text_frame
mtf.paragraphs[0].text = "Merge + sort all breakpoints → split geometry at each one"
mtf.paragraphs[0].font.size = Pt(14)
mtf.paragraphs[0].font.bold = True
mtf.paragraphs[0].font.color.rgb = WHITE
mtf.paragraphs[0].alignment = PP_ALIGN.CENTER
mtf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Result
result_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(12), Inches(0.5))
rtf = result_text.text_frame
rp = rtf.paragraphs[0]
rp.text = "Result: 245,863 segments — each stamped with county, district, area office, MPO, and RC from midpoint polygon query"
rp.font.size = Pt(15)
rp.font.bold = True
rp.font.color.rgb = ACCENT_TEAL
rp.alignment = PP_ALIGN.CENTER

# Post-split overlay note
overlay_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.7)
)
overlay_box.fill.solid()
overlay_box.fill.fore_color.rgb = LIGHT_GRAY
overlay_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
overlay_box.line.width = Pt(1)

otBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.55), Inches(10.5), Inches(0.6))
otf = otBox.text_frame
otf.word_wrap = True
op = otf.paragraphs[0]
op.text = "Post-split flagging (no further splitting):  State House / Senate / Congressional district by majority-length  |  City at \u226550% coverage"
op.font.size = Pt(12)
op.font.color.rgb = MED_GRAY
op.alignment = PP_ALIGN.CENTER


# =========================================================
# SLIDE 6: AADT Coverage Deep Dive
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "AADT Coverage Strategy"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Left column — source hierarchy
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Two Parallel GDOT-Official Sources"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = MED_BLUE

source_items = [
    ("GDOT State 2024 GDB", "Primary source — direct state traffic counts"),
    ("FHWA HPMS 2024", "Same GDOT data repackaged for federal reporting"),
    ("Combined coverage", "96.5% of segments from official sources"),
    ("Pipeline-derived fill", "Mirror / interpolation / nearest neighbor → 3.5%"),
    ("Final AADT coverage", "99.9605% — only 97 segments uncovered"),
]
for i, (label, desc) in enumerate(source_items):
    y = Inches(2.2) + Inches(i * 0.85)
    lBox = slide.shapes.add_textbox(Inches(1.0), y, Inches(5.5), Inches(0.35))
    lp = lBox.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(15)
    lp.font.bold = True
    lp.font.color.rgb = DARK_BLUE

    dBox = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.3), Inches(5.5), Inches(0.35))
    dp = dBox.text_frame.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = MED_GRAY

# Right column — HPMS historic panel
txBox3 = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "HPMS Historic AADT Panel"
p3.font.size = Pt(20)
p3.font.bold = True
p3.font.color.rgb = MED_BLUE

hpms_data = [
    ("Year", "Coverage"),
    ("2020", "94.57%"),
    ("2022", "16.17%"),
    ("2023", "16.21%"),
    ("2024", "96.31%"),
]
tbl = slide.shapes.add_table(len(hpms_data), 2, Inches(7.5), Inches(2.2), Inches(4.5), Inches(2.8)).table
tbl.columns[0].width = Inches(1.5)
tbl.columns[1].width = Inches(3.0)
for ri, (yr, cov) in enumerate(hpms_data):
    set_cell_text(tbl.cell(ri, 0), yr, font_size=13, bold=(ri == 0),
                  color=WHITE if ri == 0 else DARK_GRAY, alignment=PP_ALIGN.CENTER)
    set_cell_text(tbl.cell(ri, 1), cov, font_size=13, bold=(ri == 0),
                  color=WHITE if ri == 0 else DARK_GRAY, alignment=PP_ALIGN.CENTER)
    if ri == 0:
        style_table_header(tbl.rows[ri])
    elif ri % 2 == 0:
        for ci in range(2):
            tbl.cell(ri, ci).fill.solid()
            tbl.cell(ri, ci).fill.fore_color.rgb = TABLE_ALT_ROW

note_box = slide.shapes.add_textbox(Inches(7.5), Inches(5.2), Inches(5), Inches(1.0))
ntf = note_box.text_frame
ntf.word_wrap = True
np = ntf.paragraphs[0]
np.text = "Note: 2022/2023 lower coverage is expected — FHWA volumegroup sampling covers only federal-aid network (FC 1–5). 2021 not published by FHWA."
np.font.size = Pt(10)
np.font.color.rgb = MED_GRAY
np.font.italic = True


# =========================================================
# SLIDE 7: Future AADT Projection
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Future AADT 2044 Projection"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

fill_chain = [
    ("Step 1", "Direct GDOT official future AADT values", "Highest confidence — direct projections from state traffic data"),
    ("Step 2", "HPMS future AADT values", "Federal submission future AADT where state GDB has no projection"),
    ("Step 3", "Direction mirror", "Copy from opposite travel direction on the same route"),
    ("Step 4", "Implied growth rate (~1.17%/yr)", "Apply GDOT's observed growth rate to current AADT for remaining gaps"),
]

for i, (step, title, desc) in enumerate(fill_chain):
    y = Inches(1.6) + Inches(i * 1.2)

    # Step box
    sbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), y, Inches(1.2), Inches(0.5)
    )
    sbox.fill.solid()
    sbox.fill.fore_color.rgb = MED_BLUE
    sbox.line.fill.background()
    stf = sbox.text_frame
    stf.paragraphs[0].text = step
    stf.paragraphs[0].font.size = Pt(13)
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.color.rgb = WHITE
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow connector
    if i < 3:
        arr = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(1.45), y + Inches(0.55), Inches(0.3), Inches(0.55)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = LIGHT_BLUE
        arr.line.fill.background()

    # Title + desc
    tBox = slide.shapes.add_textbox(Inches(2.5), y - Inches(0.05), Inches(5.5), Inches(0.35))
    tp = tBox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(15)
    tp.font.bold = True
    tp.font.color.rgb = DARK_BLUE

    dBox = slide.shapes.add_textbox(Inches(2.5), y + Inches(0.28), Inches(5.5), Inches(0.35))
    dp = dBox.text_frame.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = MED_GRAY

# Result callout
result_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.5), Inches(2.0), Inches(4.2), Inches(3.5)
)
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
result_box.line.color.rgb = ACCENT_TEAL
result_box.line.width = Pt(2)

stats_future = [
    ("46,619", "segments with direct\nGDOT/HPMS forecast (19%)"),
    ("245,766", "segments with future\nAADT after fill (99.96%)"),
    ("~1.17%", "annual growth rate\nimplied by GDOT data"),
]
for i, (val, label) in enumerate(stats_future):
    y = Inches(2.2) + Inches(i * 1.1)
    vBox = slide.shapes.add_textbox(Inches(8.8), y, Inches(3.6), Inches(0.45))
    vp = vBox.text_frame.paragraphs[0]
    vp.text = val
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = ACCENT_TEAL
    vp.alignment = PP_ALIGN.CENTER

    lBox = slide.shapes.add_textbox(Inches(8.8), y + Inches(0.4), Inches(3.6), Inches(0.55))
    ltf = lBox.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(11)
    lp.font.color.rgb = MED_GRAY
    lp.alignment = PP_ALIGN.CENTER


# =========================================================
# SLIDE 8: Route Classification & Speed Limits
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Route Classification & Speed Limits"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Left — Route classification
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Signed-Route Verification"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = MED_BLUE

route_items = [
    "ROUTE_ID parsed into county, system code, route number, suffix, direction",
    "Route families: Interstate, U.S. Route, State Route, Local/Other",
    "HPMS routesigning provides 91% initial coverage (223,672 segments)",
    "GPAS reference layers are the final authority — overrides HPMS where matched",
    "99.7% agreement between state GDB and HPMS on overlap segments",
]
for i, txt in enumerate(route_items):
    y = Inches(2.2) + Inches(i * 0.7)
    bBox = slide.shapes.add_textbox(Inches(1.0), y, Inches(5.8), Inches(0.6))
    btf = bBox.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = f"\u2022  {txt}"
    bp.font.size = Pt(13)
    bp.font.color.rgb = DARK_GRAY

# Right — Speed limits
txBox3 = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Speed Limit Enrichment"
p3.font.size = Pt(20)
p3.font.bold = True
p3.font.color.rgb = MED_BLUE

speed_items = [
    ("On-System", "State highways matched by\nroute ID + milepoint overlap", "~15,000 segments"),
    ("Off-System", "Local roads matched by\nnormalized road name + county", "~30,000 segments"),
]

for i, (label, desc, count) in enumerate(speed_items):
    y = Inches(2.3) + Inches(i * 2.0)

    lbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), y, Inches(5.0), Inches(1.6)
    )
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = LIGHT_GRAY
    lbox.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    lbox.line.width = Pt(1)

    lBox = slide.shapes.add_textbox(Inches(7.8), y + Inches(0.15), Inches(4.5), Inches(0.35))
    lp = lBox.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = DARK_BLUE

    dBox = slide.shapes.add_textbox(Inches(7.8), y + Inches(0.5), Inches(3.5), Inches(0.7))
    dtf = dBox.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = MED_GRAY

    cBox = slide.shapes.add_textbox(Inches(10.5), y + Inches(0.5), Inches(1.8), Inches(0.5))
    cp = cBox.text_frame.paragraphs[0]
    cp.text = count
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = ACCENT_TEAL
    cp.alignment = PP_ALIGN.CENTER


# =========================================================
# SLIDE 9: Output Architecture
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Output Architecture"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# ETL pattern
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "ETL Pattern:  01-Raw-Data  →  02-Data-Staging  →  03-Processed-Data"
p2.font.size = Pt(16)
p2.font.color.rgb = MED_BLUE
p2.font.bold = True
p2.alignment = PP_ALIGN.CENTER

# Three output boxes
outputs = [
    ("roadway_inventory.db", "SQLite Database", "All 245,863 segments\nTabular attributes (no geometry)\n153 columns — full field lineage"),
    ("base_network.gpkg", "GeoPackage", "Segment geometry\nCounty boundary layer\nDistrict boundary layer\nSpatial-ready for GIS tools"),
    ("roadway_inventory_cleaned.csv", "Flat Table Export", "CSV for quick inspection\nSame tabular content as DB\nExcel/Pandas-friendly"),
]

for i, (fname, ftype, desc) in enumerate(outputs):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.4)

    obox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, Inches(3.8), Inches(4.2)
    )
    obox.fill.solid()
    obox.fill.fore_color.rgb = LIGHT_GRAY
    obox.line.color.rgb = MED_BLUE
    obox.line.width = Pt(1.5)

    # File icon area
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.6)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = DARK_BLUE
    icon_box.line.fill.background()
    itf = icon_box.text_frame
    itf.paragraphs[0].text = fname
    itf.paragraphs[0].font.size = Pt(13)
    itf.paragraphs[0].font.bold = True
    itf.paragraphs[0].font.color.rgb = WHITE
    itf.paragraphs[0].alignment = PP_ALIGN.CENTER
    itf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(0.4))
    tp = tBox.text_frame.paragraphs[0]
    tp.text = ftype
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = MED_BLUE
    tp.alignment = PP_ALIGN.CENTER

    dBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.7), Inches(3.2), Inches(2.2))
    dtf = dBox.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(13)
    dp.font.color.rgb = DARK_GRAY
    dp.alignment = PP_ALIGN.CENTER
    dp.line_spacing = Pt(22)


# =========================================================
# SLIDE 10: What RAPTOR Gets
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "What RAPTOR Gets"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

categories = [
    ("Category", "Key Fields"),
    ("Identity", "unique_id, ROUTE_ID, ROUTE_FAMILY, HWY_NAME, ROUTE_NUMBER"),
    ("Location", "COUNTY_CODE, COUNTY_NAME, DISTRICT, FROM/TO_MILEPOINT, geometry"),
    ("Classification", "FUNCTIONAL_CLASS, SYSTEM_CODE, ROUTE_TYPE_GDOT, SIGNED_ROUTE_FAMILY_PRIMARY"),
    ("Traffic", "AADT, TRUCK_AADT, TRUCK_PCT, K_FACTOR, D_FACTOR, VMT, FUTURE_AADT_2044"),
    ("RAPTOR-Specific", "PCT_SADT, PCT_CADT, HWY_DES"),
    ("Roadway", "NUM_LANES, SPEED_LIMIT, SURFACE_TYPE, MEDIAN_TYPE, FACILITY_TYPE, OWNERSHIP"),
    ("Pavement", "HPMS_IRI, HPMS_PSR, HPMS_RUTTING, HPMS_CRACKING_PCT"),
    ("Network", "NHS_IND, URBAN_CODE"),
    ("Data Quality", "AADT_2024_SOURCE, AADT_2024_CONFIDENCE (high/medium/low/missing)"),
]

rows, cols = len(categories), 2
col_widths = [Inches(2.0), Inches(9.5)]
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), sum(col_widths), Inches(5.5)).table
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = w

for ri, (cat, fields) in enumerate(categories):
    set_cell_text(tbl.cell(ri, 0), cat, font_size=12, bold=True,
                  color=WHITE if ri == 0 else DARK_BLUE,
                  alignment=PP_ALIGN.CENTER if ri == 0 else PP_ALIGN.LEFT)
    set_cell_text(tbl.cell(ri, 1), fields, font_size=11,
                  bold=(ri == 0),
                  color=WHITE if ri == 0 else DARK_GRAY)
    if ri == 0:
        style_table_header(tbl.rows[ri])
    elif ri % 2 == 0:
        for ci in range(cols):
            tbl.cell(ri, ci).fill.solid()
            tbl.cell(ri, ci).fill.fore_color.rgb = TABLE_ALT_ROW


# =========================================================
# SLIDE 11: Status & Next Steps
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bottom_bar(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Status & Next Steps"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

sep = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
)
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_TEAL
sep.line.fill.background()

# Status section
status_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.0)
)
status_box.fill.solid()
status_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xE8)
status_box.line.color.rgb = RGBColor(0x4C, 0xAF, 0x50)
status_box.line.width = Pt(2)

stBox = slide.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.2), Inches(0.5))
sp = stBox.text_frame.paragraphs[0]
sp.text = "Phase 1: COMPLETE"
sp.font.size = Pt(22)
sp.font.bold = True
sp.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

status_bullets = [
    "116 / 116 validation checks passing",
    "245,863 segments staged with full attribute set",
    "99.96% AADT coverage achieved",
    "All 8 data sources integrated and verified",
    "RAPTOR RoadwayData loader consuming outputs",
]
for i, txt in enumerate(status_bullets):
    y = Inches(2.3) + Inches(i * 0.4)
    bBox = slide.shapes.add_textbox(Inches(1.3), y, Inches(5.0), Inches(0.35))
    bp = bBox.text_frame.paragraphs[0]
    bp.text = f"\u2713  {txt}"
    bp.font.size = Pt(13)
    bp.font.color.rgb = DARK_GRAY

# Deferred / next section
next_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.0)
)
next_box.fill.solid()
next_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
next_box.line.color.rgb = MED_BLUE
next_box.line.width = Pt(2)

nBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(5.0), Inches(0.5))
np2 = nBox.text_frame.paragraphs[0]
np2.text = "Upcoming Phases"
np2.font.size = Pt(22)
np2.font.bold = True
np2.font.color.rgb = MED_BLUE

next_items = [
    "Phase 2: AADT station-level modeling (k-NN)",
    "Phase 3+: Additional RAPTOR scoring categories",
    "Network supplementation (optional track)",
    "Web application integration",
]
for i, txt in enumerate(next_items):
    y = Inches(2.3) + Inches(i * 0.4)
    bBox = slide.shapes.add_textbox(Inches(7.7), y, Inches(4.8), Inches(0.35))
    bp = bBox.text_frame.paragraphs[0]
    bp.text = f"\u2192  {txt}"
    bp.font.size = Pt(13)
    bp.font.color.rgb = DARK_GRAY

# Bottom callout
callout = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5))
ctf = callout.text_frame
ctf.word_wrap = True
cp = ctf.paragraphs[0]
cp.text = "The Phase 1 foundation layer is the base that all subsequent RAPTOR scoring categories build upon. With the roadway network, traffic data, and route classification verified and staged, downstream phases can proceed with a reliable statewide baseline."
cp.font.size = Pt(14)
cp.font.color.rgb = MED_GRAY
cp.font.italic = True
cp.alignment = PP_ALIGN.CENTER


# =========================================================
# Save
# =========================================================
out_path = r"d:\Jacobs\Georgia-Statewide-Data-Pipeline\00-Project-Management\tmp\Phase_1_Roadway_Pipeline_Briefing.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
